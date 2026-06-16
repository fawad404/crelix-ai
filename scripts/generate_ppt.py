"""
Crelix AI — Professional Seminar Presentation Generator
Generates: public/crelix-ai-presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy
import os

# ── Color palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x09, 0x07, 0x1C)   # near black #09071C
DARK_CARD  = RGBColor(0x14, 0x0F, 0x35)   # card bg    #140F35
PURPLE     = RGBColor(0x6E, 0x3D, 0xFF)   # brand purple
PURPLE_LT  = RGBColor(0x9B, 0x72, 0xFF)   # light purple
PINK       = RGBColor(0xFF, 0x4F, 0xCB)   # brand pink
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xBB, 0xB8, 0xD4)
GREY_DIM   = RGBColor(0x77, 0x73, 0x97)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
RED        = RGBColor(0xEF, 0x44, 0x44)

# ── Presentation setup ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank layout

# ── Low-level helpers ─────────────────────────────────────────────────────────

def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_rect_emu(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_gradient_rect(slide, x, y, w, h, color1, color2, angle=0):
    """Add a rectangle with a two-stop linear gradient via raw XML."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    # Build gradient XML
    spPr = shape.fill._xPr.getparent()
    # Remove existing fill
    for old in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill'):
        spPr.remove(old)
    for old in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill'):
        spPr.remove(old)

    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    gradFill = etree.SubElement(spPr, f'{{{ns}}}gradFill')
    gsLst    = etree.SubElement(gradFill, f'{{{ns}}}gsLst')

    gs1 = etree.SubElement(gsLst, f'{{{ns}}}gs', pos='0')
    srgb1 = etree.SubElement(gs1, f'{{{ns}}}srgbClr', val=rgb_hex(color1))

    gs2 = etree.SubElement(gsLst, f'{{{ns}}}gs', pos='100000')
    srgb2 = etree.SubElement(gs2, f'{{{ns}}}srgbClr', val=rgb_hex(color2))

    lin = etree.SubElement(gradFill, f'{{{ns}}}lin', ang=str(angle * 60000), scaled='0')
    return shape

def add_oval(slide, x, y, w, h, fill_color, transparency=0):
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name='Segoe UI', wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_multiline_text(slide, lines, x, y, w, h, font_size=14, bold=False,
                       color=WHITE, align=PP_ALIGN.LEFT, font_name='Segoe UI',
                       line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, fs, bd, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(fs)
        run.font.bold = bd
        run.font.color.rgb = clr
        run.font.name = font_name
    return txBox

def slide_bg(slide, color=DARK_BG):
    """Fill slide background."""
    add_rect(slide, 0, 0, 13.33, 7.5, color)

def accent_line(slide, x, y, w, h=0.04, color=PURPLE):
    add_rect(slide, x, y, w, h, color)

def add_pill(slide, x, y, w, h, text, fill=PURPLE, text_color=WHITE, font_size=11):
    """Rounded rectangle pill label."""
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = 'Segoe UI'
    return shape

def add_card(slide, x, y, w, h, fill=DARK_CARD):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.color.rgb = RGBColor(0x2A, 0x1E, 0x55)
    shape.line.width = Pt(1)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
def slide_title():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)

    # Big gradient blob top-right
    add_gradient_rect(slide, 6.5, -1.5, 8, 8, PURPLE, PINK, angle=45)
    # Overlay dark fade to make it subtle
    overlay = slide.shapes.add_shape(1, Inches(6.5), Inches(-1.5), Inches(8), Inches(8))
    overlay.line.fill.background()
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DARK_BG
    # Apply transparency to overlay via XML
    spPr = overlay.fill._xPr.getparent()
    solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solidFill is not None:
        srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            alpha = etree.SubElement(srgbClr, f'{{{ns}}}alpha', val='75000')  # 75% opacity

    # Small decorative purple glow circles
    add_gradient_rect(slide, 7.8, 0.8, 3.5, 3.5, PURPLE, PINK, angle=135)
    ov2 = slide.shapes.add_shape(9, Inches(7.8), Inches(0.8), Inches(3.5), Inches(3.5))
    ov2.line.fill.background()
    ov2.fill.solid()
    ov2.fill.fore_color.rgb = DARK_BG
    spPr2 = ov2.fill._xPr.getparent()
    sf2 = spPr2.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if sf2 is not None:
        sc2 = sf2.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if sc2 is not None:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            etree.SubElement(sc2, f'{{{ns}}}alpha', val='60000')

    # Dot grid decoration (small purple dots)
    for row in range(5):
        for col in range(8):
            d = slide.shapes.add_shape(9, Inches(8.8 + col*0.42), Inches(3.8 + row*0.42), Inches(0.06), Inches(0.06))
            d.line.fill.background()
            d.fill.solid()
            d.fill.fore_color.rgb = PURPLE_LT

    # Left accent bar
    add_rect(slide, 0.55, 1.8, 0.06, 3.5, PURPLE)

    # "AI MARKETING INTELLIGENCE" label pill
    add_pill(slide, 0.75, 1.9, 3.6, 0.35, "AI MARKETING INTELLIGENCE PLATFORM", PURPLE, WHITE, 10)

    # Main title — "Crelix"
    add_text(slide, "Crelix", 0.75, 2.42, 7, 1.5, font_size=82, bold=True, color=WHITE, font_name='Segoe UI Black')
    # "AI" in pink
    add_text(slide, "AI", 5.05, 2.42, 2.5, 1.5, font_size=82, bold=True, color=PINK, font_name='Segoe UI Black')

    # Tagline
    add_text(slide, "The AI that thinks like your best strategist\nand creates like your best designer.", 0.75, 4.02, 8.5, 1.0, font_size=17, bold=False, color=GREY, font_name='Segoe UI')

    # Bottom divider line
    add_rect(slide, 0.75, 5.25, 5, 0.035, PURPLE)

    # Bottom meta
    add_text(slide, "Seminar Presentation  ·  June 2026  ·  crelix.ai", 0.75, 5.38, 7, 0.4, font_size=11, color=GREY_DIM, font_name='Segoe UI')

    # Slide number indicator (bottom right)
    add_text(slide, "01", 12.3, 7.0, 0.8, 0.35, font_size=11, color=GREY_DIM, align=PP_ALIGN.RIGHT)

slide_title()

# ─────────────────────────────────────────────────────────────────────────────
# HELPER — section header template
# ─────────────────────────────────────────────────────────────────────────────
def section_header(slide, section_num, section_title, pg_num):
    # Small section number badge
    add_pill(slide, 0.55, 0.35, 0.55, 0.32, section_num, PURPLE, WHITE, 10)
    add_text(slide, section_title, 1.25, 0.28, 8, 0.5, font_size=22, bold=True, color=WHITE, font_name='Segoe UI Semibold')
    add_rect(slide, 0.55, 0.76, 12.23, 0.035, PURPLE)
    add_text(slide, str(pg_num).zfill(2), 12.3, 7.0, 0.8, 0.35, font_size=11, color=GREY_DIM, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
def slide_problem():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "01", "THE PROBLEM", 2)

    # Large background stat
    add_text(slide, "$740B", 7.5, 1.0, 5.5, 2.0, font_size=88, bold=True,
             color=RGBColor(0x1A, 0x0D, 0x44), font_name='Segoe UI Black', align=PP_ALIGN.RIGHT)
    add_text(slide, "digital advertising\nmarket — and most\nads underperform.", 7.8, 2.85, 5.0, 1.4,
             font_size=15, color=GREY_DIM, align=PP_ALIGN.RIGHT, font_name='Segoe UI')

    # Headline
    add_text(slide, "Businesses waste millions on ads that\nwere never built on real customer insight.",
             0.55, 0.95, 8.0, 1.1, font_size=24, bold=True, color=WHITE, font_name='Segoe UI Semibold')

    # 3 Pain point cards
    problems = [
        ("⏱", "Takes Days, Not Minutes",
         "Building a proper creative brief — ICP research, pain points, competitive angles — takes 3–5 business days manually."),
        ("💸", "Agency Costs Are Prohibitive",
         "A boutique creative agency charges $5,000–$25,000/month. Most SMBs and startups simply can't compete."),
        ("🎯", "Generic AI Output Doesn't Convert",
         "Current AI tools generate visually okay content with zero strategic intent — because they skip the research step entirely."),
    ]
    for i, (icon, title, desc) in enumerate(problems):
        cx = 0.55 + i * 4.26
        add_card(slide, cx, 2.32, 4.0, 2.6)
        # Top accent
        add_gradient_rect(slide, cx, 2.32, 4.0, 0.06, PURPLE, PINK)
        add_text(slide, icon, cx + 0.2, 2.45, 0.7, 0.55, font_size=24)
        add_text(slide, title, cx + 0.2, 3.05, 3.6, 0.5, font_size=14, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, desc, cx + 0.2, 3.55, 3.6, 1.1, font_size=11, color=GREY, font_name='Segoe UI')

    # Bottom stats bar
    add_rect(slide, 0.55, 5.15, 12.23, 0.8, DARK_CARD)
    stats = [("74%", "of marketers need better creative tools"),
             ("5× longer", "average time to produce one good creative brief"),
             ("10–100×", "cheaper with AI vs. traditional agencies")]
    for i, (val, lbl) in enumerate(stats):
        sx = 1.2 + i * 4.1
        add_text(slide, val, sx, 5.22, 3.5, 0.35, font_size=16, bold=True, color=PINK, font_name='Segoe UI Black')
        add_text(slide, lbl, sx, 5.56, 3.5, 0.3, font_size=10, color=GREY_DIM, font_name='Segoe UI')

slide_problem()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE SOLUTION
# ─────────────────────────────────────────────────────────────────────────────
def slide_solution():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "02", "THE SOLUTION", 3)

    # Big gradient hero area on right
    add_gradient_rect(slide, 7.0, 0.9, 6.0, 5.8, PURPLE, DARK_BG, angle=180)

    # Left content
    add_pill(slide, 0.55, 1.05, 2.2, 0.32, "INTRODUCING CRELIX AI", PURPLE, WHITE, 9)
    add_text(slide, "Intelligence First.\nCreatives Second.", 0.55, 1.52, 6.5, 1.6,
             font_size=38, bold=True, color=WHITE, font_name='Segoe UI Black')

    add_text(slide, "Crelix AI analyzes any business URL, builds a deep Marketing Brain™ using dual-AI research (Grok + Claude), then generates platform-ready ad creatives — hooks, images, benefits, CTAs — in under 3 minutes.",
             0.55, 3.22, 6.2, 1.4, font_size=13.5, color=GREY, font_name='Segoe UI')

    benefits = [
        ("✦", "No brief required — just a URL and a goal"),
        ("✦", "Real customer research, not guesswork"),
        ("✦", "Production-ready ads, not templates"),
        ("✦", "Agency quality at 1/100th the cost"),
    ]
    for i, (icon, text) in enumerate(benefits):
        add_text(slide, icon, 0.55, 4.72 + i*0.38, 0.35, 0.35, font_size=12, color=PINK, bold=True)
        add_text(slide, text, 0.9, 4.72 + i*0.38, 5.8, 0.35, font_size=13, color=WHITE, font_name='Segoe UI')

    # Right side — mock product UI card
    add_card(slide, 7.4, 1.1, 5.4, 5.55, RGBColor(0x11, 0x0A, 0x2B))
    add_gradient_rect(slide, 7.4, 1.1, 5.4, 0.06, PURPLE, PINK)
    add_text(slide, "Marketing Brain™", 7.65, 1.22, 4.9, 0.38, font_size=13, bold=True, color=WHITE, font_name='Segoe UI Semibold')
    add_text(slide, "crelix.ai/studio", 7.65, 1.6, 4.9, 0.28, font_size=10, color=GREY_DIM, font_name='Segoe UI')

    # Fake ICP display
    icp_items = [
        ("Ideal Customer Profile", "E-commerce founders, 28–42, scaling to 7-figures"),
        ("Primary Pain Point", "Spending $10K/mo on ads with declining ROAS"),
        ("Buying Trigger", "New competitor entered market, urgency to act"),
        ("Best Hook", '"Stop guessing what works. Your next winning ad starts here."'),
        ("Recommended Platform", "Meta (Primary) · Google (Secondary)"),
    ]
    for i, (label, value) in enumerate(icp_items):
        y = 2.0 + i * 0.83
        add_rect(slide, 7.6, y, 5.0, 0.72, RGBColor(0x1C, 0x12, 0x3F))
        add_text(slide, label, 7.78, y + 0.06, 4.6, 0.26, font_size=8.5, color=PURPLE_LT, bold=True, font_name='Segoe UI')
        add_text(slide, value, 7.78, y + 0.32, 4.6, 0.34, font_size=10.5, color=WHITE, font_name='Segoe UI')

slide_solution()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — HOW IT WORKS (3-step flow)
# ─────────────────────────────────────────────────────────────────────────────
def slide_how_it_works():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "03", "HOW IT WORKS", 4)

    add_text(slide, "From URL to campaign-ready creatives in 3 minutes.", 0.55, 0.88, 12, 0.5,
             font_size=19, color=GREY, font_name='Segoe UI', align=PP_ALIGN.CENTER)

    steps = [
        ("01", PURPLE, "Enter Your URL",
         "Input any business website URL and an optional campaign goal (e.g., 'increase free trial signups').",
         ["🔗 Any business URL", "🎯 Optional goal context", "⚡ Zero setup required"]),
        ("02", PINK, "Marketing Brain™ Runs",
         "Grok AI researches the business & customer. Claude synthesizes into structured intelligence with ICP, pain points, hooks & strategy.",
         ["🔍 Grok: Real-time research", "🧠 Claude: Strategic synthesis", "📊 Full ICP + angles output"]),
        ("03", GREEN, "Generate & Export",
         "Configure theme, platform, and size. Crelix generates 1–4 ad variants — full mockups with image, copy & CTA — ready to upload.",
         ["🎨 4 parallel variants", "📱 Meta / Google / TikTok", "⬇️ Export & launch"]),
    ]

    for i, (num, color, title, desc, bullets) in enumerate(steps):
        cx = 0.42 + i * 4.3
        # Card
        add_card(slide, cx, 1.55, 4.05, 5.5)
        add_rect(slide, cx, 1.55, 4.05, 0.07, color)

        # Step number circle
        add_oval(slide, cx + 0.18, 1.72, 0.7, 0.7, color)
        add_text(slide, num, cx + 0.18, 1.72, 0.7, 0.7, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Segoe UI Black')

        # Arrow between steps
        if i < 2:
            add_text(slide, "→", cx + 4.05, 3.8, 0.3, 0.5, font_size=22, color=GREY_DIM, bold=True)

        add_text(slide, title, cx + 0.18, 2.52, 3.7, 0.5, font_size=16, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, desc, cx + 0.18, 3.06, 3.68, 1.05, font_size=11, color=GREY, font_name='Segoe UI')

        for j, bullet in enumerate(bullets):
            by = 4.2 + j * 0.5
            add_rect(slide, cx + 0.18, by + 0.11, 0.03, 0.22, color)
            add_text(slide, bullet, cx + 0.28, by, 3.55, 0.44, font_size=11.5, color=WHITE, font_name='Segoe UI')

    # Bottom time badge
    add_gradient_rect(slide, 4.5, 7.0, 4.33, 0.38, PURPLE, PINK)
    add_text(slide, "⚡  Total time: under 3 minutes  ⚡", 4.5, 7.0, 4.33, 0.38,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

slide_how_it_works()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — MARKETING BRAIN
# ─────────────────────────────────────────────────────────────────────────────
def slide_marketing_brain():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "04", "MARKETING BRAIN™ — THE INTELLIGENCE LAYER", 5)

    # Left column — description
    add_text(slide, "The world's first pre-generation\nmarketing intelligence engine.", 0.55, 1.0, 5.5, 1.1,
             font_size=24, bold=True, color=WHITE, font_name='Segoe UI Black')
    add_text(slide, "Before generating a single image or word of copy, Crelix AI runs a deep dual-AI research pipeline — powered by Grok's real-time internet intelligence and Claude's structured synthesis — to build a complete dossier on your business, your customer, and your market.",
             0.55, 2.18, 5.3, 1.6, font_size=12, color=GREY, font_name='Segoe UI')

    # Dual AI pill
    add_gradient_rect(slide, 0.55, 3.9, 5.3, 0.55, PURPLE, PINK)
    add_text(slide, "🔬 Grok AI  →  Research & Intel       🧩 Claude AI  →  Synthesis & Output",
             0.55, 3.9, 5.3, 0.55, font_size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

    details = [
        "Scrapes and analyzes business website content",
        "Identifies exact business type and positioning",
        "Maps ICP demographics and psychographics",
        "Extracts top 5 customer pain points",
        "Surfaces 3 buying triggers and purchase objections",
        "Recommends ad angles, platforms, and offers",
        "Generates 5 hooks, 3 headlines, 3 CTAs",
    ]
    for i, d in enumerate(details):
        add_text(slide, "✓", 0.55, 4.6 + i*0.35, 0.3, 0.33, font_size=11, color=GREEN, bold=True)
        add_text(slide, d, 0.85, 4.6 + i*0.35, 5.0, 0.33, font_size=11, color=GREY, font_name='Segoe UI')

    # Right — output cards grid
    outputs = [
        (PURPLE, "ICP Profile", "Who is buying, why, in what life situation"),
        (PINK,   "Pain Points", "5 real customer fears in their own words"),
        (GREEN,  "Buying Triggers", "Events that create urgency to buy NOW"),
        (AMBER,  "Ad Angles", "3 proven creative angles with reasoning"),
        (PURPLE_LT, "Hook Copy", "5 scroll-stopping opening lines, ready to use"),
        (RGBColor(0x22,0xD3,0xEE), "Platform Strategy", "Which platform, why, and in what priority"),
    ]
    for i, (clr, title, sub) in enumerate(outputs):
        col = i % 2
        row = i // 2
        cx = 6.25 + col * 3.35
        cy = 1.0 + row * 2.12
        add_card(slide, cx, cy, 3.15, 1.95)
        add_rect(slide, cx, cy, 3.15, 0.055, clr)
        add_text(slide, title, cx + 0.18, cy + 0.14, 2.8, 0.38, font_size=13, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, sub, cx + 0.18, cy + 0.56, 2.8, 1.1, font_size=11, color=GREY, font_name='Segoe UI')

slide_marketing_brain()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — CREATIVE STUDIO
# ─────────────────────────────────────────────────────────────────────────────
def slide_creative_studio():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "05", "CREATIVE STUDIO™ — THE EXECUTION LAYER", 6)

    add_text(slide, "From intelligence to production-ready ad creatives\nin one click.", 0.55, 0.95, 8.5, 0.85,
             font_size=22, bold=True, color=WHITE, font_name='Segoe UI Black')

    # Left feature list
    features = [
        ("🎨", "10 Visual Themes", "Lighting styles from Dark Premium to Ocean Blue — affects mood, not subject matter"),
        ("📱", "3 Platform Modes", "Meta (square), Google (horizontal), TikTok (vertical) — composition adapts automatically"),
        ("⚡", "4 Parallel Variants", "Up to 4 ad variants generated simultaneously using different hooks and angles"),
        ("🖼", "Full Ad Mockup", "Brand bar + Hook text + 3 benefit bullets + CTA button overlaid on the image"),
        ("🧠", "ICP-Derived Images", "Characters and scenes come from the Marketing Brain ICP — not random AI guesses"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        cy = 2.0 + i * 0.95
        add_rect(slide, 0.55, cy, 5.5, 0.84, DARK_CARD)
        add_text(slide, icon, 0.72, cy + 0.08, 0.55, 0.65, font_size=20)
        add_text(slide, title, 1.32, cy + 0.06, 4.5, 0.32, font_size=13, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, desc, 1.32, cy + 0.38, 4.5, 0.4, font_size=10.5, color=GREY, font_name='Segoe UI')

    # Right — ad mockup preview (drawn with shapes)
    mx, my = 6.55, 0.95
    mw, mh = 6.4, 6.25
    add_card(slide, mx, my, mw, mh, RGBColor(0x18, 0x10, 0x38))
    add_text(slide, "LIVE AD MOCKUP PREVIEW", mx, my + 0.08, mw, 0.28,
             font_size=8, color=GREY_DIM, align=PP_ALIGN.CENTER, font_name='Segoe UI')

    # Facebook-style top bar
    add_rect(slide, mx + 0.12, my + 0.42, mw - 0.24, 0.52, RGBColor(0x22, 0x14, 0x4A))
    add_oval(slide, mx + 0.18, my + 0.48, 0.4, 0.4, PURPLE)
    add_text(slide, "C", mx + 0.18, my + 0.48, 0.4, 0.4, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Your Brand Name", mx + 0.65, my + 0.48, 3.0, 0.22, font_size=11, bold=True, color=WHITE, font_name='Segoe UI Semibold')
    add_text(slide, "Sponsored ·", mx + 0.65, my + 0.7, 2.0, 0.2, font_size=8.5, color=GREY_DIM, font_name='Segoe UI')

    # Image area (simulated with gradient)
    add_gradient_rect(slide, mx + 0.12, my + 0.98, mw - 0.24, 2.9, RGBColor(0x2A,0x18,0x66), DARK_BG, angle=180)
    # Image placeholder text
    add_text(slide, "[ AI-Generated Background Image ]", mx + 0.12, my + 2.1, mw - 0.24, 0.5,
             font_size=11, color=RGBColor(0x3A,0x28,0x70), align=PP_ALIGN.CENTER, italic=True)

    # Overlay content
    add_text(slide, "⬤ YourBrand", mx + 0.22, my + 1.1, 2.5, 0.28, font_size=9, bold=True, color=WHITE)
    add_text(slide, '"Stop wasting ad spend on\ncreatives that don\'t convert."', mx + 0.22, my + 1.5, mw - 0.4, 0.65,
             font_size=14, bold=True, color=WHITE, font_name='Segoe UI Black')
    for j, b in enumerate(["✓ ICP-targeted creative", "✓ Platform-optimized size", "✓ Conversion-tested copy"]):
        add_text(slide, b, mx + 0.22, my + 2.25 + j*0.26, 4.2, 0.24, font_size=9.5, color=RGBColor(0xCC,0xBB,0xFF))
    add_pill(slide, mx + 0.22, my + 3.1, 1.85, 0.36, "Get Started Free →", PURPLE, WHITE, 10)

    # Footer card
    add_rect(slide, mx + 0.12, my + 3.95, mw - 0.24, 0.7, RGBColor(0x1A, 0x10, 0x3A))
    add_text(slide, "crelix.ai", mx + 0.22, my + 4.0, 3.0, 0.24, font_size=9, color=GREY_DIM)
    add_text(slide, "Crelix AI — AI Marketing Platform", mx + 0.22, my + 4.24, 4.5, 0.24, font_size=10, bold=True, color=WHITE)
    add_pill(slide, mx + 4.8, my + 4.1, 1.4, 0.32, "Learn More", RGBColor(0x1E,0x4A,0xFF), WHITE, 9)

    add_text(slide, "↕ 4 variants generated in parallel", mx, my + 4.82, mw, 0.3,
             font_size=9.5, color=GREY_DIM, align=PP_ALIGN.CENTER, italic=True)

slide_creative_studio()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — TECHNOLOGY ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
def slide_technology():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "06", "TECHNOLOGY ARCHITECTURE", 7)

    add_text(slide, "Dual-AI Pipeline — the only platform using both the best research model\nand the best synthesis model, in a single seamless workflow.",
             0.55, 0.92, 12.2, 0.75, font_size=15, color=GREY, font_name='Segoe UI', align=PP_ALIGN.CENTER)

    # Pipeline flow — horizontal
    nodes = [
        (PURPLE, "STEP 1", "Website\nScraper", "Fetches title, H1, meta,\nclean content from URL\n\npython fetch + HTML parsing\nServer-side only"),
        (RGBColor(0xF5,0x9E,0x0B), "STEP 2", "Grok AI\nResearch", "Real-time business\n& customer intelligence\n\ngrok-4 series model\nMax 1200 tokens output"),
        (PINK, "STEP 3", "Claude AI\nSynthesis", "Structured JSON output\nvia tool_use constraint\n\nclaude-sonnet-4-6\nGuaranteed schema"),
        (GREEN, "STEP 4", "Image\nGeneration", "4 parallel prompts\nwith ICP-derived\n\ngrok-imagine-image-quality\nPromise.allSettled()"),
        (RGBColor(0x22,0xD3,0xEE), "STEP 5", "Ad\nCompositor", "Brand bar + overlay\nhooks + benefits + CTA\n\nReact CSS overlay\nFull mockup output"),
    ]
    for i, (clr, label, title, detail) in enumerate(nodes):
        cx = 0.5 + i * 2.55
        # Arrow
        if i > 0:
            add_text(slide, "→", cx - 0.35, 2.2, 0.35, 0.4, font_size=18, color=GREY_DIM, bold=True)
        # Node card
        add_card(slide, cx, 1.78, 2.3, 3.8)
        add_gradient_rect(slide, cx, 1.78, 2.3, 0.07, clr, clr)
        add_pill(slide, cx + 0.12, 1.92, 1.0, 0.26, label, clr, WHITE, 8)
        add_text(slide, title, cx + 0.12, 2.28, 2.05, 0.65, font_size=14, bold=True, color=WHITE, font_name='Segoe UI Black')
        add_text(slide, detail, cx + 0.12, 3.0, 2.06, 2.4, font_size=9.5, color=GREY, font_name='Segoe UI')

    # Bottom tech stack bar
    add_rect(slide, 0.5, 5.8, 12.3, 1.32, DARK_CARD)
    add_text(slide, "TECH STACK", 0.7, 5.88, 2.0, 0.3, font_size=9, bold=True, color=PURPLE_LT, font_name='Segoe UI')
    stack = [
        ("Next.js 16", "App Router + Turbopack"),
        ("React 19", "Client components + hooks"),
        ("xAI / Grok API", "Research + image generation"),
        ("Anthropic Claude", "claude-sonnet-4-6 via tool_use"),
        ("Vercel Edge", "Serverless, maxDuration 300s"),
    ]
    for i, (tech, desc) in enumerate(stack):
        sx = 0.7 + i * 2.45
        add_text(slide, tech, sx, 6.22, 2.3, 0.3, font_size=11.5, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, desc, sx, 6.54, 2.3, 0.5, font_size=9.5, color=GREY_DIM, font_name='Segoe UI')

slide_technology()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — MARKET OPPORTUNITY
# ─────────────────────────────────────────────────────────────────────────────
def slide_market():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "07", "MARKET OPPORTUNITY", 8)

    add_text(slide, "Positioned at the intersection of the two fastest-growing\nsoftware categories of the decade.", 0.55, 0.9, 12.2, 0.75,
             font_size=17, color=GREY, align=PP_ALIGN.CENTER)

    # TAM / SAM / SOM concentric rings (simulated with nested rectangles)
    # TAM
    add_oval(slide, 0.9, 1.85, 5.8, 5.2, RGBColor(0x1E, 0x12, 0x45))
    add_text(slide, "TAM", 2.5, 2.1, 2.5, 0.35, font_size=11, bold=True, color=GREY_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "$740B", 1.5, 2.45, 4.5, 0.55, font_size=28, bold=True, color=RGBColor(0x4A,0x30,0x99), align=PP_ALIGN.CENTER, font_name='Segoe UI Black')
    add_text(slide, "Global Digital\nAdvertising", 2.0, 2.98, 3.5, 0.55, font_size=10, color=GREY_DIM, align=PP_ALIGN.CENTER)

    # SAM
    add_oval(slide, 1.55, 2.65, 4.5, 3.65, RGBColor(0x2A, 0x18, 0x60))
    add_text(slide, "SAM", 3.0, 3.05, 1.8, 0.3, font_size=10, bold=True, color=GREY_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "$107B", 2.3, 3.38, 3.2, 0.5, font_size=24, bold=True, color=PURPLE_LT, align=PP_ALIGN.CENTER, font_name='Segoe UI Black')
    add_text(slide, "AI in Marketing\nby 2028", 2.5, 3.88, 2.8, 0.45, font_size=9.5, color=GREY_DIM, align=PP_ALIGN.CENTER)

    # SOM
    add_oval(slide, 2.35, 3.55, 2.9, 2.05, RGBColor(0x3A, 0x22, 0x80))
    add_text(slide, "SOM", 3.4, 3.78, 0.95, 0.28, font_size=9, bold=True, color=PINK, align=PP_ALIGN.CENTER)
    add_text(slide, "$28B", 2.65, 4.08, 2.3, 0.42, font_size=20, bold=True, color=PINK, align=PP_ALIGN.CENTER, font_name='Segoe UI Black')
    add_text(slide, "AI Creative\nTools 2028", 2.8, 4.5, 2.0, 0.38, font_size=9, color=GREY_DIM, align=PP_ALIGN.CENTER)

    # Right side — market stats
    stats = [
        (PURPLE,    "26.1%", "CAGR", "AI in Marketing 2021–2028"),
        (PINK,      "$8B",   "NOW",  "AI Creative Tools market today"),
        (GREEN,     "14,000+","TOOLS","Marketing technology solutions globally"),
        (AMBER,     "0%",    "GAP",  "Platforms with intelligence-first workflow"),
    ]
    for i, (clr, val, badge, label) in enumerate(stats):
        cy = 1.1 + i * 1.55
        add_card(slide, 7.0, cy, 5.8, 1.35)
        add_rect(slide, 7.0, cy, 0.06, 1.35, clr)
        add_pill(slide, 7.18, cy + 0.08, 0.9, 0.28, badge, clr, WHITE, 8)
        add_text(slide, val, 7.18, cy + 0.42, 2.2, 0.7, font_size=34, bold=True, color=WHITE, font_name='Segoe UI Black')
        add_text(slide, label, 9.2, cy + 0.5, 3.4, 0.55, font_size=12.5, color=GREY, font_name='Segoe UI')

slide_market()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — COMPETITIVE LANDSCAPE
# ─────────────────────────────────────────────────────────────────────────────
def slide_competitive():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "08", "COMPETITIVE LANDSCAPE", 9)

    add_text(slide, "No competitor combines intelligence research + image generation + strategic copy\nin a single workflow. That gap is Crelix's entire market position.",
             0.55, 0.88, 12.2, 0.65, font_size=13.5, color=GREY, align=PP_ALIGN.CENTER)

    # Table headers
    headers = ["Platform", "Marketing\nBrain", "AI Image\nGen", "Strategic\nCopy", "ICP\nResearch", "Price/mo", "Our Edge"]
    col_w   = [2.2, 1.38, 1.38, 1.38, 1.38, 1.2, 2.7]
    col_x   = [0.42]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w + 0.05)

    # Header row
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_gradient_rect(slide, x, 1.72, w, 0.52, PURPLE, RGBColor(0x4A,0x1A,0xBB))
        add_text(slide, h, x + 0.05, 1.72, w, 0.52, font_size=9.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

    # Rows
    rows = [
        ("Crelix AI", "✅ Full", "✅ Grok", "✅ Claude", "✅ Yes", "$49–499", "Only intelligence-first platform"),
        ("AdCreative.ai", "❌ None", "✅ SD", "⚠️ Basic", "❌ No", "$29–149", "No research, template-only"),
        ("Jasper AI", "❌ None", "❌ None", "✅ Strong", "❌ No", "$49–125", "Writing only, no images"),
        ("Canva AI", "❌ None", "✅ Design", "⚠️ Weak", "❌ No", "$15–55", "Designer tool, not advertiser"),
        ("Copy.ai", "❌ None", "❌ None", "✅ Good", "❌ No", "$49–249", "Copy only, no strategy"),
    ]
    for i, row in enumerate(rows):
        is_crelix = i == 0
        row_bg = RGBColor(0x1C, 0x10, 0x42) if is_crelix else (DARK_CARD if i % 2 == 0 else RGBColor(0x11, 0x0A, 0x2A))
        for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
            add_rect(slide, x, 2.28 + i * 0.74, w, 0.68, row_bg)
            clr = WHITE if (j == 0 or is_crelix) else GREY
            if cell.startswith("✅"):
                clr = GREEN if not is_crelix else WHITE
            elif cell.startswith("❌"):
                clr = RED
            elif cell.startswith("⚠️"):
                clr = AMBER
            if is_crelix and j == 0:
                add_gradient_rect(slide, x, 2.28 + i * 0.74, w, 0.68, RGBColor(0x3A,0x15,0x99), PURPLE)
            fs = 11 if j == 0 else 11
            if j == 6:
                fs = 9.5
            add_text(slide, cell, x + 0.06, 2.28 + i * 0.74 + 0.12, w - 0.1, 0.44, font_size=fs,
                     bold=(j == 0 and is_crelix), color=clr, align=PP_ALIGN.CENTER, font_name='Segoe UI')

    # Callout bottom
    add_gradient_rect(slide, 0.42, 6.0, 12.48, 0.55, PURPLE, PINK)
    add_text(slide, "💡  Key Insight: Every competitor assumes you already know your customer. Crelix is the only platform that finds out for you first.",
             0.6, 6.0, 12.2, 0.55, font_size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

slide_competitive()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — BUSINESS MODEL
# ─────────────────────────────────────────────────────────────────────────────
def slide_business_model():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "09", "BUSINESS MODEL", 10)

    add_text(slide, "Subscription SaaS — three tiers aligned to customer segment value.",
             0.55, 0.9, 12.2, 0.42, font_size=17, color=GREY, align=PP_ALIGN.CENTER)

    plans = [
        (PURPLE, "STARTER", "$49", "/month", "For independent business owners and solopreneurs running their own ads.",
         ["10 Marketing Brains / month", "20 ad creatives / month", "2 platform modes", "All 10 visual themes", "Email support"]),
        (PINK, "PRO", "$149", "/month", "For in-house marketing teams needing volume and campaign strategy tools.",
         ["Unlimited Marketing Brains", "100 creatives / month", "All 3 platforms", "Campaign Builder access", "Creative Psychology Report", "Priority support"]),
        (GREEN, "AGENCY", "$499", "/month", "For agencies managing multiple client brands simultaneously.",
         ["Unlimited everything", "25 brand profiles (Brand Memory)", "White-label exports", "Client management dashboard", "Direct Meta Ads integration", "Dedicated success manager"]),
    ]
    for i, (clr, name, price, period, desc, features) in enumerate(plans):
        cx = 0.55 + i * 4.26
        # Card with slightly different height for Pro (recommended)
        extra = 0.2 if i == 1 else 0
        cy = 1.48 - extra
        ch = 5.62 + extra

        if i == 1:
            # Glow effect for recommended
            add_rect(slide, cx - 0.05, cy - 0.05, 4.1, ch + 0.1, PURPLE)
        add_card(slide, cx, cy, 4.0, ch, DARK_CARD)
        add_gradient_rect(slide, cx, cy, 4.0, 0.07, clr, clr)

        if i == 1:
            add_pill(slide, cx + 1.2, cy - 0.32, 1.6, 0.28, "MOST POPULAR", PINK, WHITE, 9)

        add_text(slide, name, cx + 0.2, cy + 0.18, 3.6, 0.38, font_size=11, bold=True,
                 color=clr, font_name='Segoe UI Black', align=PP_ALIGN.CENTER)
        add_text(slide, price, cx + 0.2, cy + 0.55, 3.6, 0.82, font_size=48, bold=True,
                 color=WHITE, font_name='Segoe UI Black', align=PP_ALIGN.CENTER)
        add_text(slide, period, cx + 0.2, cy + 1.35, 3.6, 0.3, font_size=12, color=GREY_DIM,
                 align=PP_ALIGN.CENTER)
        add_rect(slide, cx + 0.2, cy + 1.7, 3.6, 0.03, RGBColor(0x2A,0x1E,0x55))
        add_text(slide, desc, cx + 0.2, cy + 1.82, 3.6, 0.7, font_size=10.5, color=GREY,
                 font_name='Segoe UI', align=PP_ALIGN.CENTER)
        for j, feat in enumerate(features):
            add_text(slide, "✓", cx + 0.22, cy + 2.62 + j*0.46, 0.3, 0.38, font_size=11, color=clr, bold=True)
            add_text(slide, feat, cx + 0.52, cy + 2.62 + j*0.46, 3.28, 0.38, font_size=11, color=WHITE, font_name='Segoe UI')

    # Footer note
    add_text(slide, "Annual plans save 20%  ·  Free trial: 3 Marketing Brains + 5 creatives, no credit card required",
             0.55, 7.1, 12.2, 0.3, font_size=10, color=GREY_DIM, align=PP_ALIGN.CENTER, italic=True)

slide_business_model()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — FINANCIAL PROJECTIONS
# ─────────────────────────────────────────────────────────────────────────────
def slide_financials():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "10", "FINANCIAL PROJECTIONS", 11)

    add_text(slide, "Conservative 3-year growth model — bootstrapped Year 1, seed-funded Years 2–3.",
             0.55, 0.9, 12.2, 0.4, font_size=15, color=GREY, align=PP_ALIGN.CENTER)

    # Left — big ARR numbers
    years = [
        (PURPLE, "Year 1", "$465K", "ARR", "~480 customers", "$80 avg/mo"),
        (PINK,   "Year 2", "$1.33M", "ARR", "~1,350 customers", "$100 avg/mo"),
        (GREEN,  "Year 3", "$3.0M", "ARR", "~3,200 customers", "$115 avg/mo"),
    ]
    for i, (clr, yr, arr, lbl, cust, arpu) in enumerate(years):
        cy = 1.45 + i * 1.8
        add_card(slide, 0.55, cy, 5.8, 1.62)
        add_rect(slide, 0.55, cy, 0.07, 1.62, clr)
        add_text(slide, yr, 0.78, cy + 0.1, 1.5, 0.32, font_size=11, bold=True, color=clr, font_name='Segoe UI Black')
        add_text(slide, arr, 0.78, cy + 0.38, 2.5, 0.75, font_size=38, bold=True, color=WHITE, font_name='Segoe UI Black')
        add_text(slide, lbl, 3.2, cy + 0.52, 1.0, 0.35, font_size=11, color=GREY_DIM)
        add_text(slide, cust, 0.78, cy + 1.12, 3.0, 0.32, font_size=12, color=GREY, font_name='Segoe UI')
        add_text(slide, arpu, 3.8, cy + 1.12, 2.3, 0.32, font_size=11, color=clr, bold=True, font_name='Segoe UI')

    # Right — key metrics table
    add_text(slide, "SaaS Health Metrics", 6.7, 1.45, 6.3, 0.38, font_size=14, bold=True, color=WHITE, font_name='Segoe UI Semibold')
    metrics = [
        ("Gross Margin",         "85%",    "88%",    "90%"),
        ("Monthly Churn",        "5.0%",   "3.5%",   "2.5%"),
        ("LTV : CAC Ratio",      "28×",    "36×",    "46×"),
        ("Net Revenue Retention","105%",   "115%",   "125%"),
        ("Break-even MRR",       "$12K",   "✅ Met",  "✅ Met"),
        ("ARR Growth",           "—",      "+185%",  "+126%"),
    ]
    mh = ["Metric", "Year 1", "Year 2", "Year 3"]
    mw = [2.85, 1.08, 1.08, 1.08]
    mx_arr = [6.65, 9.5, 10.58, 11.66]

    for j, (h, x, w) in enumerate(zip(mh, mx_arr, mw)):
        add_gradient_rect(slide, x, 1.9, w, 0.42, PURPLE, RGBColor(0x4A,0x1A,0xBB))
        add_text(slide, h, x + 0.04, 1.9, w, 0.42, font_size=9.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

    row_colors = [DARK_CARD, RGBColor(0x11,0x0A,0x2A)]
    for i, (label, y1, y2, y3) in enumerate(metrics):
        bg = row_colors[i % 2]
        cells = [label, y1, y2, y3]
        for j, (cell, x, w) in enumerate(zip(cells, mx_arr, mw)):
            add_rect(slide, x, 2.36 + i * 0.58, w, 0.55, bg)
            clr = WHITE if j == 0 else GREY
            if "✅" in cell:
                clr = GREEN
            elif "%" in cell and j > 0:
                clr = PURPLE_LT
            add_text(slide, cell, x + 0.04, 2.36 + i * 0.58 + 0.08, w - 0.06, 0.38,
                     font_size=10.5, color=clr, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                     bold=(j == 0), font_name='Segoe UI')

    # Seed round callout
    add_gradient_rect(slide, 6.65, 5.92, 6.1, 1.1, RGBColor(0x1E,0x10,0x55), PURPLE)
    add_text(slide, "Raising Pre-Seed", 6.85, 6.0, 5.7, 0.3, font_size=11, bold=True, color=WHITE, font_name='Segoe UI Semibold')
    add_text(slide, "$300K–500K", 6.85, 6.28, 5.7, 0.45, font_size=22, bold=True, color=PINK, font_name='Segoe UI Black')
    add_text(slide, "To accelerate product roadmap, first hire & GTM", 6.85, 6.72, 5.7, 0.28,
             font_size=10, color=GREY, font_name='Segoe UI')

slide_financials()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — SWOT ANALYSIS
# ─────────────────────────────────────────────────────────────────────────────
def slide_swot():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "11", "SWOT ANALYSIS", 12)

    quadrants = [
        (PURPLE,  "STRENGTHS",     "S", 0.42, 1.05, [
            "Only intelligence-first AI ad platform",
            "Dual-AI (Grok + Claude) architecture",
            "10–100× cheaper than agencies",
            "Under 3 minutes from URL to full creative set",
            "Platform-specific optimization (Meta/Google/TikTok)",
            "ICP-derived image generation (not random scenes)",
        ]),
        (AMBER,   "WEAKNESSES",    "W", 6.95, 1.05, [
            "Early-stage brand — limited public case studies",
            "Dependency on 3rd-party AI APIs (xAI, Anthropic)",
            "No direct ad platform integrations yet",
            "Small founding team — bandwidth constraints",
            "Brand Memory not yet shipped",
        ]),
        (GREEN,   "OPPORTUNITIES", "O", 0.42, 4.35, [
            "$107B AI in marketing market growing 26% CAGR",
            "Agency channel — 1 partner = 10–50 client brands",
            "AdCreative.ai users actively seeking alternatives",
            "International markets virtually untapped",
            "Platform integrations create deep switching cost moat",
        ]),
        (RED,     "THREATS",       "T", 6.95, 4.35, [
            "Meta/Google building native AI creative tools",
            "AI API cost increases reducing margin",
            "Well-funded competitors adding intelligence layers",
            "Regulatory changes for AI-generated ads",
            "Brand skepticism about AI creative quality",
        ]),
    ]
    for clr, title, letter, cx, cy, items in quadrants:
        add_card(slide, cx, cy, 6.0, 2.98, DARK_CARD)
        add_gradient_rect(slide, cx, cy, 6.0, 0.07, clr, clr)

        # Large letter watermark
        add_text(slide, letter, cx + 4.5, cy + 0.2, 1.4, 1.4, font_size=72, bold=True,
                 color=RGBColor(max(0, clr[0]-80), max(0, clr[1]-80), max(0, clr[2]-80)),
                 font_name='Segoe UI Black')

        add_text(slide, title, cx + 0.18, cy + 0.14, 4.5, 0.38, font_size=13, bold=True,
                 color=clr, font_name='Segoe UI Black')
        for j, item in enumerate(items):
            add_text(slide, "·  " + item, cx + 0.18, cy + 0.6 + j * 0.39, 5.5, 0.37,
                     font_size=10.5, color=WHITE, font_name='Segoe UI')

    # Center cross dividers
    add_rect(slide, 6.48, 1.05, 0.04, 6.28, RGBColor(0x2A,0x1E,0x55))
    add_rect(slide, 0.42, 4.0, 12.51, 0.04, RGBColor(0x2A,0x1E,0x55))

slide_swot()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — ROADMAP & MILESTONES
# ─────────────────────────────────────────────────────────────────────────────
def slide_roadmap():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "12", "ROADMAP & MILESTONES", 13)

    add_text(slide, "12 months to platform dominance in AI marketing intelligence.",
             0.55, 0.88, 12.2, 0.38, font_size=15, color=GREY, align=PP_ALIGN.CENTER)

    # Timeline line
    add_gradient_rect(slide, 0.55, 4.0, 12.2, 0.07, PURPLE, PINK)

    milestones = [
        (PURPLE, "Q2 2026", "Launch", ["Marketing Brain™ live", "Creative Studio™ live", "50 beta users onboarded"]),
        (PINK,   "Q3 2026", "Traction", ["First 100 paying customers", "$10K MRR milestone", "Campaign Builder ships"]),
        (AMBER,  "Q4 2026", "Growth", ["Brand Memory (Supabase)", "20 Agency partners", "$500K ARR target"]),
        (GREEN,  "Q1 2027", "Scale", ["Meta Ads integration", "Seed round closes", "Head of Growth hired"]),
        (RGBColor(0x22,0xD3,0xEE), "Q4 2027", "Expand", ["$2M ARR milestone", "Google Ads integration", "Series A preparation"]),
    ]

    for i, (clr, quarter, phase, items) in enumerate(milestones):
        cx = 0.55 + i * 2.46

        # Top card (above timeline)
        add_card(slide, cx, 1.12, 2.3, 2.74)
        add_rect(slide, cx, 1.12, 2.3, 0.06, clr)
        add_pill(slide, cx + 0.08, 1.22, 1.4, 0.28, quarter, clr, WHITE, 9)
        add_text(slide, phase, cx + 0.12, 1.6, 2.05, 0.38, font_size=14, bold=True, color=WHITE, font_name='Segoe UI Black')
        for j, item in enumerate(items):
            add_text(slide, "→ " + item, cx + 0.12, 2.08 + j * 0.55, 2.05, 0.48, font_size=10, color=GREY, font_name='Segoe UI')

        # Circle on timeline
        add_oval(slide, cx + 0.75, 3.74, 0.55, 0.55, clr)
        add_text(slide, str(i+1), cx + 0.75, 3.74, 0.55, 0.55, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Bottom KPI
        kpi_labels = ["50 users", "100 cust.", "500K ARR", "Seed ✓", "2M ARR"]
        add_rect(slide, cx, 4.45, 2.3, 0.5, DARK_CARD)
        add_text(slide, kpi_labels[i], cx + 0.1, 4.5, 2.1, 0.38, font_size=12, bold=True, color=clr, align=PP_ALIGN.CENTER, font_name='Segoe UI Black')

    # Bottom note
    add_text(slide, "Series A target: Q2 2028 · $8–12M · $40–60M valuation  ·  Exit horizon: acquisition or IPO at $20M+ ARR",
             0.55, 5.1, 12.2, 0.35, font_size=10.5, color=GREY_DIM, align=PP_ALIGN.CENTER, italic=True)

slide_roadmap()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — GO-TO-MARKET STRATEGY
# ─────────────────────────────────────────────────────────────────────────────
def slide_gtm():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "13", "GO-TO-MARKET STRATEGY", 14)

    add_text(slide, "Product-led growth with community amplification — let the output sell itself.",
             0.55, 0.88, 12.2, 0.38, font_size=15, color=GREY, align=PP_ALIGN.CENTER)

    # PLG flywheel diagram (3 circular stages)
    flywheel = [
        (PURPLE, "ACQUIRE", "Free trial: 3 Marketing Brains + 5 creatives\nNo credit card. Instant value delivery.\nOrganic SEO + community presence"),
        (PINK,   "ACTIVATE", "User runs first Marketing Brain analysis\nSees their ICP + pain points + hooks\nAha moment: 'This is exactly right'"),
        (GREEN,  "EXPAND",   "Upgrades to paid plan\nRefers colleagues + agencies\nCase study → social proof → more acquires"),
    ]
    for i, (clr, stage, desc) in enumerate(flywheel):
        cx = 0.55 + i * 4.26
        # Big circle
        add_oval(slide, cx + 0.3, 1.35, 3.45, 3.45, RGBColor(max(0,clr[0]-120), max(0,clr[1]-120), max(0,clr[2]-120)))
        add_oval(slide, cx + 0.55, 1.6, 2.95, 2.95, RGBColor(max(0,clr[0]-60), max(0,clr[1]-60), max(0,clr[2]-60)))
        add_text(slide, stage, cx + 0.3, 2.05, 3.45, 0.5, font_size=18, bold=True, color=clr, align=PP_ALIGN.CENTER, font_name='Segoe UI Black')
        add_text(slide, desc, cx + 0.3, 2.62, 3.45, 1.6, font_size=10.5, color=WHITE, align=PP_ALIGN.CENTER, font_name='Segoe UI')

        if i < 2:
            add_text(slide, "⟶", cx + 3.95, 2.85, 0.55, 0.55, font_size=22, color=GREY_DIM, bold=True)

    # Channels grid
    channels = [
        ("📝", "Content SEO", "High-intent blog content targeting performance marketers. Ranks for 'AI ad creative generator', 'marketing intelligence platform'."),
        ("🐦", "Twitter / LinkedIn", "Founder-led distribution. Daily marketing intelligence teardowns, product demos, behind-the-scenes AI building journey."),
        ("🤝", "Agency Partnerships", "20 agency partners in Year 1. Each agency = 10–50 client brands at Agency tier pricing."),
        ("🎯", "Free Audit Funnel", "Any business gets partial Marketing Brain free — ICP + top 3 pain points — generating product-qualified leads."),
    ]
    for i, (icon, title, desc) in enumerate(channels):
        col = i % 2
        row = i // 2
        cx = 0.55 + col * 6.2
        cy = 5.0 + row * 1.15
        add_card(slide, cx, cy, 6.0, 1.05)
        add_text(slide, icon, cx + 0.15, cy + 0.08, 0.55, 0.75, font_size=22)
        add_text(slide, title, cx + 0.75, cy + 0.08, 2.5, 0.32, font_size=12, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, desc, cx + 0.75, cy + 0.42, 5.1, 0.55, font_size=9.5, color=GREY, font_name='Segoe UI')

slide_gtm()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — THE TEAM
# ─────────────────────────────────────────────────────────────────────────────
def slide_team():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "14", "THE TEAM", 15)

    # Decorative gradient right side
    add_gradient_rect(slide, 7.5, 0.85, 5.5, 6.5, PURPLE, DARK_BG, angle=180)

    # Founder card - left
    add_card(slide, 0.55, 1.05, 6.7, 4.2)
    add_gradient_rect(slide, 0.55, 1.05, 6.7, 0.07, PURPLE, PINK)

    # Avatar circle
    add_oval(slide, 0.78, 1.28, 1.35, 1.35, PURPLE)
    add_text(slide, "FA", 0.78, 1.28, 1.35, 1.35, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Segoe UI Black')

    add_text(slide, "Fawad Ansari", 2.35, 1.3, 4.6, 0.5, font_size=22, bold=True, color=WHITE, font_name='Segoe UI Black')
    add_text(slide, "Founder & CEO · Chief Product Officer", 2.35, 1.8, 4.6, 0.35, font_size=12, color=PINK, font_name='Segoe UI')
    add_text(slide, "crelix.ai", 2.35, 2.15, 4.6, 0.3, font_size=11, color=GREY_DIM, font_name='Segoe UI')

    add_rect(slide, 0.78, 2.56, 6.28, 0.03, RGBColor(0x2A,0x1E,0x55))

    bio_points = [
        "Built the complete dual-AI Marketing Brain pipeline from concept to production",
        "Designed the Creative Studio™ with ICP-derived image generation architecture",
        "Full-stack engineering across Next.js 16, React 19, xAI Grok, Anthropic Claude",
        "Deep expertise in digital marketing, AI systems, and SaaS product development",
        "Platform built and validated with real users — live, functional, shipping features",
    ]
    for j, pt in enumerate(bio_points):
        add_text(slide, "◆", 0.78, 2.72 + j*0.44, 0.3, 0.38, font_size=10, color=PURPLE_LT, bold=True)
        add_text(slide, pt, 1.1, 2.72 + j*0.44, 5.92, 0.38, font_size=11.5, color=WHITE, font_name='Segoe UI')

    # Hiring plan card
    add_card(slide, 0.55, 5.38, 6.7, 1.72)
    add_text(slide, "IMMEDIATE HIRING PLAN", 0.75, 5.46, 6.0, 0.32, font_size=10, bold=True, color=PURPLE_LT, font_name='Segoe UI')
    hires = [
        (PINK,   "Hire #1 — Head of Growth", "PLG SaaS acquisition, content marketing, agency partnerships (Q3 2026)"),
        (GREEN,  "Hire #2 — AI Engineer", "Pipeline optimization, roadmap acceleration, inference scaling (Q4 2026)"),
    ]
    for k, (clr, role, detail) in enumerate(hires):
        add_rect(slide, 0.75, 5.82 + k*0.6, 0.06, 0.44, clr)
        add_text(slide, role, 0.88, 5.82 + k*0.6, 3.0, 0.24, font_size=11, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, detail, 0.88, 6.06 + k*0.6, 6.0, 0.24, font_size=10, color=GREY, font_name='Segoe UI')

    # Right side — why us
    add_text(slide, "Why This Team\nWins", 8.0, 1.2, 4.6, 0.9, font_size=28, bold=True, color=WHITE, font_name='Segoe UI Black')
    advantages = [
        ("🧠", "Built-in domain expertise", "Founder has lived the problem — wasted budget on generic ads with no strategic foundation."),
        ("⚡", "Technical speed advantage", "Sole technical founder means zero coordination lag. Architecture decisions happen in minutes, not meetings."),
        ("🎯", "Product-market intuition", "Every feature exists because a real customer need existed first — no padding, no bloat."),
        ("🌐", "Network & visibility", "Active in performance marketing communities with growing audience of target customers."),
    ]
    for i, (icon, title, desc) in enumerate(advantages):
        cy = 2.28 + i * 1.22
        add_rect(slide, 7.85, cy, 4.9, 1.08, DARK_CARD)
        add_text(slide, icon, 7.98, cy + 0.1, 0.55, 0.75, font_size=22)
        add_text(slide, title, 8.55, cy + 0.1, 4.1, 0.3, font_size=12, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, desc, 8.55, cy + 0.44, 4.1, 0.55, font_size=10, color=GREY, font_name='Segoe UI')

slide_team()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — WHY NOW
# ─────────────────────────────────────────────────────────────────────────────
def slide_why_now():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    section_header(slide, "15", "WHY NOW", 16)

    add_text(slide, "Three converging forces have created a narrow window\nto define this category — and Crelix is moving first.",
             0.55, 0.88, 12.2, 0.75, font_size=20, bold=True, color=WHITE, font_name='Segoe UI Black', align=PP_ALIGN.CENTER)

    trends = [
        (PURPLE, "01", "AI Models Reached\nProduction Quality",
         "In 2023, AI-generated images were novelties. In 2026, they are indistinguishable from photography. Grok Aurora and Claude operate at a level where AI output is not a compromise — it's the preferred format for speed-to-market brands.",
         ["Grok Aurora: photorealistic ad images", "Claude: strategist-quality synthesis", "Inference costs fell 90% in 18 months"]),
        (PINK, "02", "The 'Intelligence Gap'\nIs Now Visible",
         "Advertisers who adopted first-gen AI tools (Jasper, AdCreative) are now reporting performance plateaus. They got the execution but not the strategy. The market has learned it needs both — and is actively searching for a solution.",
         ["AdCreative.ai users report 'generic feel'", "Marketers understand ICP matters for AI", "Demand for 'smarter' AI tools is explicit"]),
        (GREEN, "03", "Digital Ad Spend\nIs at an Inflection",
         "With third-party cookies gone, first-party data scarce, and CPMs rising, creative quality is now the primary variable in ad performance. The businesses that win will be those with the most strategic, most tested, most targeted creative — which is exactly what Crelix delivers.",
         ["Creative is the new targeting", "Ad testing volume = performance edge", "SMBs need agency-quality at SMB prices"]),
    ]
    for i, (clr, num, title, body, bullets) in enumerate(trends):
        cx = 0.42 + i * 4.3
        add_card(slide, cx, 1.82, 4.05, 5.22)
        add_gradient_rect(slide, cx, 1.82, 4.05, 0.08, clr, clr)

        add_text(slide, num, cx + 0.18, 2.0, 0.65, 0.55, font_size=28, bold=True,
                 color=RGBColor(max(0,clr[0]-60), max(0,clr[1]-60), max(0,clr[2]-60)), font_name='Segoe UI Black')
        add_text(slide, title, cx + 0.18, 2.62, 3.7, 0.65, font_size=15, bold=True, color=WHITE, font_name='Segoe UI Black')
        add_text(slide, body, cx + 0.18, 3.35, 3.7, 1.55, font_size=10.5, color=GREY, font_name='Segoe UI')
        for j, b in enumerate(bullets):
            add_text(slide, "✦ " + b, cx + 0.18, 5.0 + j*0.34, 3.7, 0.32, font_size=10, color=clr, font_name='Segoe UI')

slide_why_now()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — THANK YOU / CTA
# ─────────────────────────────────────────────────────────────────────────────
def slide_thankyou():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)

    # Full gradient background glow
    add_gradient_rect(slide, 2.5, -0.5, 10, 6, PURPLE, PINK, angle=45)
    # Dark overlay
    ov = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    ov.line.fill.background()
    ov.fill.solid()
    ov.fill.fore_color.rgb = DARK_BG
    spPr = ov.fill._xPr.getparent()
    sf = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if sf is not None:
        sc = sf.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if sc is not None:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            etree.SubElement(sc, f'{{{ns}}}alpha', val='82000')

    # Dot grid
    for row in range(8):
        for col in range(15):
            d = slide.shapes.add_shape(9, Inches(0.3 + col*0.88), Inches(0.3 + row*0.9), Inches(0.055), Inches(0.055))
            d.line.fill.background()
            d.fill.solid()
            d.fill.fore_color.rgb = RGBColor(0x2A, 0x18, 0x55)

    # Center content
    add_text(slide, "Thank You", 1.5, 1.2, 10.33, 1.5, font_size=78, bold=True,
             color=WHITE, font_name='Segoe UI Black', align=PP_ALIGN.CENTER)

    add_gradient_rect(slide, 4.2, 2.85, 4.93, 0.07, PURPLE, PINK)

    add_text(slide, "Let's build the future of AI marketing together.",
             1.5, 3.05, 10.33, 0.5, font_size=18, color=GREY, align=PP_ALIGN.CENTER, font_name='Segoe UI')

    # CTA boxes
    ctas = [
        (PURPLE, "🌐  crelix.ai", "Live Platform"),
        (PINK,   "📧  fawadanxari31@gmail.com", "Get in Touch"),
        (GREEN,  "🚀  Try Free — No Credit Card", "3 Brains + 5 Creatives"),
    ]
    for i, (clr, label, sub) in enumerate(ctas):
        cx = 1.55 + i * 3.5
        add_card(slide, cx, 3.78, 3.25, 1.2, DARK_CARD)
        add_rect(slide, cx, 3.78, 3.25, 0.055, clr)
        add_text(slide, label, cx + 0.15, 3.9, 2.95, 0.42, font_size=13, bold=True, color=WHITE, font_name='Segoe UI Semibold')
        add_text(slide, sub, cx + 0.15, 4.32, 2.95, 0.28, font_size=10, color=GREY_DIM, font_name='Segoe UI')

    # Bottom tagline
    add_text(slide, "\"Intelligence First. Creatives Second. Results Always.\"",
             1.5, 5.3, 10.33, 0.45, font_size=14, color=PURPLE_LT, align=PP_ALIGN.CENTER, italic=True)

    add_text(slide, "Crelix AI  ·  AI Marketing Intelligence Platform  ·  2026",
             1.5, 6.85, 10.33, 0.38, font_size=10, color=GREY_DIM, align=PP_ALIGN.CENTER)

slide_thankyou()

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), '..', 'public', 'crelix-ai-presentation.pptx')
out_path = os.path.abspath(out_path)
prs.save(out_path)
print(f"✅  Saved: {out_path}")
print(f"    Slides: {len(prs.slides)}")
